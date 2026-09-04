from __future__ import annotations

import io
from html import escape
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from .office_renderer import has_libreoffice, render_office_to_pdf


def _fallback_powerpoint_to_pdf(source: Path, output: Path) -> str:
    presentation = Presentation(source)
    styles = getSampleStyleSheet()
    story = []

    for slide_index, slide in enumerate(presentation.slides):
        if slide_index:
            story.append(PageBreak())
        story.append(Paragraph(f'Slide {slide_index + 1}', styles['Heading2']))
        story.append(Spacer(1, 10))

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = Image(io.BytesIO(shape.image.blob))
                    image._restrictSize(6.5 * inch, 7.5 * inch)
                    story.append(image)
                    story.append(Spacer(1, 8))
                except Exception:
                    pass
                continue

            if getattr(shape, 'has_table', False):
                data = []
                for row in shape.table.rows:
                    data.append([
                        Paragraph(escape((cell.text or '').strip()).replace('\n', '<br/>'), styles['Normal'])
                        for cell in row.cells
                    ])
                if data:
                    table = Table(data, hAlign='LEFT')
                    table.setStyle(TableStyle([
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                    ]))
                    story.append(table)
                    story.append(Spacer(1, 8))
                continue

            text = ''
            if getattr(shape, 'has_text_frame', False):
                text = '\n'.join(
                    (''.join(run.text for run in paragraph.runs) or paragraph.text or '').strip()
                    for paragraph in shape.text_frame.paragraphs
                ).strip()
            elif hasattr(shape, 'text'):
                text = (shape.text or '').strip()

            if text:
                story.append(Paragraph(escape(text).replace('\n', '<br/>'), styles['Normal']))
                story.append(Spacer(1, 6))

    pdf = SimpleDocTemplate(str(output), pagesize=letter, leftMargin=54, rightMargin=54, topMargin=54, bottomMargin=54)
    pdf.build(story or [Paragraph('Converted presentation is empty.', styles['Normal'])])
    return str(output)


def powerpoint_to_pdf(input_path: str | Path, output_path: str | Path) -> str:
    source = Path(input_path)
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    if has_libreoffice():
        rendered = render_office_to_pdf(source, output)
        if rendered:
            return rendered
        rendered = render_office_to_pdf(source, output)
        if rendered:
            return rendered
        raise RuntimeError(
            'LibreOffice is installed but could not convert this PowerPoint file. '
            'Close other LibreOffice windows and try again.'
        )

    return _fallback_powerpoint_to_pdf(source, output)
