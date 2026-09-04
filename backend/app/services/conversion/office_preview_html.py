from __future__ import annotations

import base64
import csv
import io
from html import escape
from pathlib import Path

from docx import Document
from docx.document import Document as DocxDocument
from docx.oxml.ns import qn
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _iter_docx_blocks(document: DocxDocument):
    body = document.element.body
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            yield DocxParagraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield DocxTable(child, document)


def _docx_run_images(paragraph: DocxParagraph) -> list[tuple[str, bytes]]:
    images: list[tuple[str, bytes]] = []
    for run in paragraph.runs:
        drawings = run._element.xpath('.//a:blip')
        for blip in drawings:
            embed = blip.get(qn('r:embed'))
            if not embed:
                continue
            try:
                part = paragraph.part.related_parts[embed]
            except KeyError:
                continue
            content_type = getattr(part, 'content_type', 'image/png') or 'image/png'
            images.append((content_type, part.blob))
    return images


def docx_to_preview_html(source: Path) -> str:
    document = Document(source)
    chunks: list[str] = ['<article class="office-preview office-preview-docx">']

    for block in _iter_docx_blocks(document):
        if isinstance(block, DocxParagraph):
            for content_type, blob in _docx_run_images(block):
                encoded = base64.b64encode(blob).decode('ascii')
                chunks.append(
                    f'<p class="office-preview-image"><img src="data:{content_type};base64,{encoded}" alt="" /></p>'
                )
            text = (block.text or '').strip()
            if text:
                style = (block.style.name if block.style is not None else '') or ''
                tag = 'h2' if style.startswith('Heading') else 'p'
                chunks.append(f'<{tag}>{escape(text)}</{tag}>')
        elif isinstance(block, DocxTable):
            chunks.append('<table>')
            for row in block.rows:
                chunks.append('<tr>')
                for cell in row.cells:
                    chunks.append(f'<td>{escape((cell.text or "").strip())}</td>')
                chunks.append('</tr>')
            chunks.append('</table>')

    chunks.append('</article>')
    return '\n'.join(chunks)


def excel_to_preview_html(source: Path) -> str:
    if source.suffix.lower() == '.csv':
        text = source.read_text(encoding='utf-8', errors='replace')
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)
        chunks = ['<article class="office-preview office-preview-excel"><table>']
        for row in rows[:200]:
            chunks.append('<tr>' + ''.join(f'<td>{escape(str(cell))}</td>' for cell in row) + '</tr>')
        chunks.append('</table></article>')
        return '\n'.join(chunks)

    workbook = load_workbook(source, read_only=True, data_only=True)
    chunks = ['<article class="office-preview office-preview-excel">']
    for sheet in workbook.worksheets:
        chunks.append(f'<h2>{escape(sheet.title)}</h2><table>')
        row_count = 0
        for row in sheet.iter_rows(values_only=True):
            row_count += 1
            if row_count > 200:
                chunks.append('<tr><td colspan="99">Preview truncated after 200 rows.</td></tr>')
                break
            cells = ['' if value is None else str(value) for value in row]
            chunks.append('<tr>' + ''.join(f'<td>{escape(cell)}</td>' for cell in cells) + '</tr>')
        chunks.append('</table>')
    workbook.close()
    chunks.append('</article>')
    return '\n'.join(chunks)


def powerpoint_to_preview_html(source: Path) -> str:
    presentation = Presentation(source)
    chunks = ['<article class="office-preview office-preview-pptx">']
    for index, slide in enumerate(presentation.slides, start=1):
        chunks.append(f'<section class="office-preview-slide"><h2>Slide {index}</h2>')
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                encoded = base64.b64encode(image.blob).decode('ascii')
                content_type = image.content_type or 'image/png'
                chunks.append(
                    f'<p class="office-preview-image"><img src="data:{content_type};base64,{encoded}" alt="" /></p>'
                )
            elif getattr(shape, 'has_text_frame', False):
                for paragraph in shape.text_frame.paragraphs:
                    text = ''.join(run.text for run in paragraph.runs).strip() or (paragraph.text or '').strip()
                    if text:
                        chunks.append(f'<p>{escape(text)}</p>')
            elif hasattr(shape, 'text') and shape.text and shape.text.strip():
                chunks.append(f'<p>{escape(shape.text.strip())}</p>')
            if getattr(shape, 'has_table', False):
                chunks.append('<table>')
                for row in shape.table.rows:
                    chunks.append('<tr>')
                    for cell in row.cells:
                        chunks.append(f'<td>{escape((cell.text or "").strip())}</td>')
                    chunks.append('</tr>')
                chunks.append('</table>')
        chunks.append('</section>')
    chunks.append('</article>')
    return '\n'.join(chunks)


def office_to_preview_html(source: Path) -> str:
    suffix = source.suffix.lower()
    if suffix in {'.doc', '.docx'}:
        return docx_to_preview_html(source)
    if suffix in {'.xls', '.xlsx', '.csv'}:
        return excel_to_preview_html(source)
    if suffix in {'.ppt', '.pptx'}:
        return powerpoint_to_preview_html(source)
    raise ValueError(f'Unsupported preview type: {suffix}')
