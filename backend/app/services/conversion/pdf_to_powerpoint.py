from __future__ import annotations

import io
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from backend.app.services.pdf.structure import page_reading_items


def _emu(points: float) -> int:
    return int(points * 12700)


def pdf_to_powerpoint(input_path: str | Path, output_path: str | Path) -> str:
    """
    Convert PDF → editable PowerPoint.
    Each page becomes a slide with positioned text boxes, tables-as-text, and embedded images.
    Full-page screenshot slides are not used.
    """
    source = Path(input_path)
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]

    with fitz.open(str(source)) as document:
        for page in document:
            page_width = page.rect.width
            page_height = page.rect.height
            presentation.slide_width = Inches(max(page_width / 72.0, 1))
            presentation.slide_height = Inches(max(page_height / 72.0, 1))
            slide = presentation.slides.add_slide(blank_layout)

            items = page_reading_items(document, page)
            if not items:
                text = (page.get_text('text') or '').strip() or '[Empty page]'
                box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(max(page_width / 72.0 - 1, 1)), Inches(1.5))
                box.text_frame.word_wrap = True
                box.text_frame.text = text
                continue

            for kind, payload in items:
                if kind == 'table':
                    # Represent tables as an editable text grid on the slide.
                    left = max(payload.bbox[0], 8)
                    top = max(payload.bbox[1], 8)
                    width = max(payload.bbox[2] - payload.bbox[0], 80)
                    height = max(payload.bbox[3] - payload.bbox[1], 40)
                    box = slide.shapes.add_textbox(_emu(left), _emu(top), _emu(width), _emu(height))
                    frame = box.text_frame
                    frame.word_wrap = True
                    frame.clear()
                    for row_index, row in enumerate(payload.rows):
                        paragraph = frame.paragraphs[0] if row_index == 0 else frame.add_paragraph()
                        paragraph.text = ' | '.join(row)
                        paragraph.font.size = Pt(12)
                        paragraph.font.bold = row_index == 0
                    continue

                block = payload
                left = max(block.bbox[0], 4)
                top = max(block.bbox[1], 4)
                width = max(block.bbox[2] - block.bbox[0], 24)
                height = max(block.bbox[3] - block.bbox[1], 18)

                if block.kind == 'image' and block.image_bytes:
                    stream = io.BytesIO(block.image_bytes)
                    slide.shapes.add_picture(stream, _emu(left), _emu(top), width=_emu(width), height=_emu(height))
                    continue

                if block.kind != 'text':
                    continue

                box = slide.shapes.add_textbox(_emu(left), _emu(top), _emu(width), _emu(max(height, 20)))
                frame = box.text_frame
                frame.word_wrap = True
                # Reset default paragraph, then rebuild styled runs.
                frame.text = ''
                first = True
                for line in block.lines:
                    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
                    first = False
                    paragraph.alignment = PP_ALIGN.LEFT
                    for span in line.spans:
                        run = paragraph.add_run()
                        run.text = span.text
                        run.font.size = Pt(max(8, min(span.size, 40)))
                        run.font.bold = span.bold
                        run.font.italic = span.italic
                        run.font.color.rgb = RGBColor(*span.color)

    if not presentation.slides:
        presentation.slide_width = Inches(10)
        presentation.slide_height = Inches(7.5)
        presentation.slides.add_slide(blank_layout)

    presentation.save(output)
    return str(output)
